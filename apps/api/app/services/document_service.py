"""
MiniDoc - Document Generation Service
Generate PDFs, Excel, Word, PowerPoint documents.
"""
import io
from typing import List, Dict, Any, Optional
from datetime import datetime
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
import openpyxl
from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
from openpyxl.utils import get_column_letter
from docx import Document as DocxDocument
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class DocumentService:
    """Service for generating various document types."""

    # ==================== PDF Generation ====================

    async def generate_pdf(
        self,
        title: str,
        content: str,
        template: str = "default",
        metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Generate a PDF document.
        
        Args:
            title: Document title
            content: Main content (markdown-like text)
            template: Template style (default, report, invoice, letter)
            metadata: Additional metadata (author, subject, etc.)
            
        Returns:
            Dict with 'pdf_bytes' and 'page_count'
        """
        try:
            buffer = io.BytesIO()
            
            # Create document
            doc = SimpleDocTemplate(
                buffer,
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            # Styles
            styles = getSampleStyleSheet()
            
            # Custom styles
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                alignment=TA_CENTER,
                spaceAfter=30
            )
            
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontSize=14,
                spaceBefore=20,
                spaceAfter=10
            )
            
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['Normal'],
                fontSize=11,
                alignment=TA_JUSTIFY,
                spaceBefore=6,
                spaceAfter=6
            )
            
            # Build content
            story = []
            
            # Title
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 12))
            
            # Metadata
            if metadata:
                meta_text = f"Author: {metadata.get('author', 'Unknown')}"
                if metadata.get('date'):
                    meta_text += f" | Date: {metadata['date']}"
                story.append(Paragraph(meta_text, styles['Normal']))
                story.append(Spacer(1, 24))
            
            # Parse and add content
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                para = para.strip()
                if not para:
                    continue
                
                # Check for headings (lines starting with #)
                if para.startswith('# '):
                    story.append(Paragraph(para[2:], heading_style))
                elif para.startswith('## '):
                    story.append(Paragraph(para[3:], heading_style))
                else:
                    # Regular paragraph
                    story.append(Paragraph(para, body_style))
            
            # Build PDF
            doc.build(story)
            buffer.seek(0)
            pdf_bytes = buffer.read()
            
            return {
                "success": True,
                "pdf_bytes": pdf_bytes,
                "page_count": doc.page  # Approximate
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "pdf_bytes": None
            }

    async def generate_pdf_report(
        self,
        title: str,
        sections: List[Dict[str, str]],
        metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Generate a structured PDF report with sections.
        
        Args:
            title: Report title
            sections: List of {heading, content} dictionaries
            metadata: Report metadata
            
        Returns:
            Dict with 'pdf_bytes'
        """
        try:
            buffer = io.BytesIO()
            
            doc = SimpleDocTemplate(
                buffer,
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            styles = getSampleStyleSheet()
            
            title_style = ParagraphStyle(
                'ReportTitle',
                parent=styles['Heading1'],
                fontSize=24,
                alignment=TA_CENTER,
                spaceAfter=30
            )
            
            heading_style = ParagraphStyle(
                'SectionHeading',
                parent=styles['Heading2'],
                fontSize=14,
                spaceBefore=20,
                spaceAfter=10,
                textColor=colors.darkblue
            )
            
            body_style = ParagraphStyle(
                'ReportBody',
                parent=styles['Normal'],
                fontSize=11,
                alignment=TA_JUSTIFY,
                spaceBefore=6,
                spaceAfter=6
            )
            
            story = []
            
            # Title page
            story.append(Spacer(1, 100))
            story.append(Paragraph(title, title_style))
            
            if metadata:
                story.append(Spacer(1, 30))
                if metadata.get('subtitle'):
                    story.append(Paragraph(metadata['subtitle'], styles['Normal']))
                if metadata.get('author'):
                    story.append(Paragraph(f"Prepared by: {metadata['author']}", styles['Normal']))
                if metadata.get('date'):
                    story.append(Paragraph(f"Date: {metadata['date']}", styles['Normal']))
            
            story.append(PageBreak())
            
            # Table of contents placeholder
            story.append(Paragraph("Table of Contents", heading_style))
            for idx, section in enumerate(sections, 1):
                story.append(Paragraph(f"{idx}. {section.get('heading', 'Section')}", body_style))
            story.append(PageBreak())
            
            # Sections
            for idx, section in enumerate(sections, 1):
                heading = section.get('heading', f'Section {idx}')
                content = section.get('content', '')
                
                story.append(Paragraph(heading, heading_style))
                story.append(Paragraph(content, body_style))
                story.append(Spacer(1, 20))
            
            doc.build(story)
            buffer.seek(0)
            
            return {
                "success": True,
                "pdf_bytes": buffer.read()
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e)
            }

    # ==================== Excel Generation ====================

    async def generate_spreadsheet(
        self,
        title: str,
        data: List[Dict[str, Any]],
        headers: List[str] = None,
        sheet_name: str = "Sheet1"
    ) -> Dict[str, Any]:
        """
        Generate an Excel spreadsheet.
        
        Args:
            title: Spreadsheet title (filename)
            data: List of dictionaries (rows)
            headers: Column headers (if None, uses keys from first dict)
            sheet_name: Name of the sheet
            
        Returns:
            Dict with 'xlsx_bytes'
        """
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet_name
            
            # Get headers from data if not provided
            if not headers and data:
                headers = list(data[0].keys())
            
            if not headers:
                headers = []
            
            # Style definitions
            header_font = Font(bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_alignment = Alignment(horizontal="center", vertical="center")
            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Write headers
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment
                cell.border = thin_border
            
            # Write data
            for row_idx, row_data in enumerate(data, 2):
                for col_idx, header in enumerate(headers, 1):
                    value = row_data.get(header, "")
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.border = thin_border
                    cell.alignment = Alignment(vertical="center")
            
            # Auto-adjust column widths
            for col_idx, header in enumerate(headers, 1):
                max_length = len(str(header))
                for row in ws.iter_rows(min_row=2, min_col=col_idx, max_col=col_idx):
                    for cell in row:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                ws.column_dimensions[get_column_letter(col_idx)].width = min(max_length + 2, 50)
            
            # Save to bytes
            buffer = io.BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            
            return {
                "success": True,
                "xlsx_bytes": buffer.read(),
                "row_count": len(data),
                "column_count": len(headers)
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "xlsx_bytes": None
            }

    async def generate_spreadsheet_from_table(
        self,
        title: str,
        table_data: List[List[str]],
        has_header: bool = True
    ) -> Dict[str, Any]:
        """
        Generate an Excel spreadsheet from a 2D array.
        
        Args:
            title: Spreadsheet title
            table_data: 2D array of values
            has_header: Whether first row is header
            
        Returns:
            Dict with 'xlsx_bytes'
        """
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Data"
            
            header_font = Font(bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            
            for row_idx, row in enumerate(table_data, 1):
                for col_idx, value in enumerate(row, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    
                    if has_header and row_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill
            
            # Auto-adjust columns
            for col in ws.columns:
                max_length = 0
                column = col[0].column_letter
                for cell in col:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                ws.column_dimensions[column].width = max_length + 2
            
            buffer = io.BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            
            return {
                "success": True,
                "xlsx_bytes": buffer.read(),
                "row_count": len(table_data),
                "column_count": len(table_data[0]) if table_data else 0
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e)
            }

    # ==================== Word Document Generation ====================

    async def generate_docx(
        self,
        title: str,
        content: str,
        formatting: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Generate a Word document.
        
        Args:
            title: Document title
            content: Main content (text with newlines for paragraphs)
            formatting: Optional formatting options
            
        Returns:
            Dict with 'docx_bytes'
        """
        try:
            doc = DocxDocument()
            
            # Title
            title_para = doc.add_heading(title, level=0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Add content
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                para = para.strip()
                if para:
                    if para.startswith('# '):
                        doc.add_heading(para[2:], level=1)
                    elif para.startswith('## '):
                        doc.add_heading(para[3:], level=2)
                    elif para.startswith('### '):
                        doc.add_heading(para[4:], level=3)
                    else:
                        doc.add_paragraph(para)
            
            # Save to bytes
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            
            return {
                "success": True,
                "docx_bytes": buffer.read(),
                "paragraph_count": len(doc.paragraphs)
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "docx_bytes": None
            }

    async def generate_docx_report(
        self,
        title: str,
        sections: List[Dict[str, str]],
        metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Generate a Word document with sections.
        
        Args:
            title: Document title
            sections: List of {heading, content}
            metadata: Document metadata
            
        Returns:
            Dict with 'docx_bytes'
        """
        try:
            doc = DocxDocument()
            
            # Title
            doc.add_heading(title, level=0)
            
            # Metadata
            if metadata:
                meta_para = doc.add_paragraph()
                if metadata.get('author'):
                    meta_para.add_run(f"Author: {metadata['author']}\n")
                if metadata.get('date'):
                    meta_para.add_run(f"Date: {metadata['date']}\n")
                doc.add_paragraph()  # Spacer
            
            # Sections
            for section in sections:
                heading = section.get('heading', '')
                content = section.get('content', '')
                
                if heading:
                    doc.add_heading(heading, level=1)
                if content:
                    doc.add_paragraph(content)
            
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            
            return {
                "success": True,
                "docx_bytes": buffer.read()
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e)
            }

    # ==================== PowerPoint Generation ====================

    async def generate_pptx(
        self,
        title: str,
        slides: List[Dict[str, str]],
        theme: str = "default"
    ) -> Dict[str, Any]:
        """
        Generate a PowerPoint presentation.
        
        Args:
            title: Presentation title
            slides: List of {title, content} for each slide
            theme: Theme style
            
        Returns:
            Dict with 'pptx_bytes'
        """
        try:
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title
            subtitle_shape.text = f"Generated by MiniDoc\n{datetime.now().strftime('%B %d, %Y')}"
            
            # Content slides
            bullet_slide_layout = prs.slide_layouts[1]
            
            for slide_data in slides:
                slide_title = slide_data.get('title', '')
                slide_content = slide_data.get('content', '')
                
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                # Title
                shapes = slide.shapes
                title_shape = shapes.title
                title_shape.text = slide_title
                
                # Content
                body_shape = shapes.placeholders[1]
                text_frame = body_shape.text_frame
                
                # Split content into bullet points
                points = slide_content.split('\n')
                for idx, point in enumerate(points):
                    point = point.strip()
                    if not point:
                        continue
                    
                    if idx == 0:
                        text_frame.text = point
                    else:
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0
            
            # Save to bytes
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            return {
                "success": True,
                "pptx_bytes": buffer.read(),
                "slide_count": len(prs.slides)
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "pptx_bytes": None
            }

    async def generate_pptx_from_outline(
        self,
        title: str,
        outline: str
    ) -> Dict[str, Any]:
        """
        Generate a PowerPoint presentation from a text outline.
        
        Args:
            title: Presentation title
            outline: Outline with slide headers and bullet points
            
        Returns:
            Dict with 'pptx_bytes'
        """
        try:
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            
            # Parse outline
            lines = outline.strip().split('\n')
            current_slide = None
            current_content = []
            slides_data = []
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # Slide title (starts with # or number.)
                if line.startswith('#') or (line[0].isdigit() and '.' in line[:3]):
                    if current_slide:
                        slides_data.append({
                            'title': current_slide,
                            'content': '\n'.join(current_content)
                        })
                    current_slide = line.lstrip('#').lstrip('0123456789.').strip()
                    current_content = []
                else:
                    # Bullet point
                    current_content.append(line.lstrip('-*').strip())
            
            # Add last slide
            if current_slide:
                slides_data.append({
                    'title': current_slide,
                    'content': '\n'.join(current_content)
                })
            
            # Create slides
            bullet_slide_layout = prs.slide_layouts[1]
            for slide_data in slides_data:
                slide = prs.slides.add_slide(bullet_slide_layout)
                slide.shapes.title.text = slide_data['title']
                
                body_shape = slide.shapes.placeholders[1]
                text_frame = body_shape.text_frame
                
                points = slide_data['content'].split('\n')
                for idx, point in enumerate(points):
                    if idx == 0:
                        text_frame.text = point
                    else:
                        p = text_frame.add_paragraph()
                        p.text = point
            
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            return {
                "success": True,
                "pptx_bytes": buffer.read(),
                "slide_count": len(prs.slides)
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e)
            }


# Singleton instance
document_service = DocumentService()
